"""Eager document → text extraction for project files.

Runs in the backend process at upload time so the model receives plain text
directly instead of having to open the file with a tool. Dispatch is by file
extension. Each extractor imports its heavy dependency lazily so a missing
package degrades to ``("", "failed")`` rather than crashing the upload — the
chat path then falls back to letting the model read the raw file itself.
"""

from __future__ import annotations

import csv
import io
import json
from pathlib import Path
from typing import Literal

ExtractStatus = Literal["ok", "failed", "skipped"]

# Hard cap per file so one huge spreadsheet can't blow the context budget.
# The project-context assembler applies a second, global cap across all files.
_MAX_CHARS = 200_000

# Extensions we read as UTF-8 text verbatim (no library needed).
_PLAIN_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".yaml", ".yml",
    ".xml", ".html", ".htm", ".log", ".rst", ".ini", ".toml", ".text",
}


def _truncate(text: str) -> str:
    if len(text) <= _MAX_CHARS:
        return text
    return text[:_MAX_CHARS] + f"\n\n[...truncated: showing first {_MAX_CHARS} characters]"


def extract_text(data: bytes, filename: str, mime_type: str = "") -> tuple[str, ExtractStatus]:
    """Extract plain text from a document's raw bytes.

    Returns ``(text, "ok")`` on success, ``("", "failed")`` when the format is
    known but extraction errored or its library is missing, and
    ``("", "skipped")`` for formats we deliberately don't extract (images,
    archives, binaries) — those stay file-only.
    """
    ext = Path(filename).suffix.lower()

    try:
        if ext in _PLAIN_TEXT_EXTS:
            return _truncate(_decode(data)), "ok"
        if ext == ".docx":
            return _truncate(_extract_docx(data)), "ok"
        if ext == ".xlsx":
            return _truncate(_extract_xlsx(data)), "ok"
        if ext == ".pdf":
            return _truncate(_extract_pdf(data)), "ok"
        if ext == ".pptx":
            return _truncate(_extract_pptx(data)), "ok"
    except Exception:  # noqa: BLE001 — any extractor failure → fall back to file-only
        return "", "failed"

    # Unknown / binary (.png, .zip, .doc, …) — leave for the model to open.
    return "", "skipped"


def _decode(data: bytes) -> str:
    """Best-effort UTF-8 decode; CSV/TSV get a light normalisation pass."""
    return data.decode("utf-8", errors="replace")


def _extract_docx(data: bytes) -> str:
    from docx import Document  # python-docx

    doc = Document(io.BytesIO(data))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    out = io.StringIO()
    for ws in wb.worksheets:
        out.write(f"# Sheet: {ws.title}\n")
        writer = csv.writer(out)
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                writer.writerow(["" if c is None else str(c) for c in row])
        out.write("\n")
    wb.close()
    return out.getvalue().strip()


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n\n".join(p.strip() for p in pages if p.strip())
    if not text.strip():
        # Scanned/image PDF — no text layer. Let the model decide (OCR via tool).
        raise ValueError("no extractable text layer")
    return text


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation  # python-pptx

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"# Slide {idx}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
    return "\n".join(parts)
